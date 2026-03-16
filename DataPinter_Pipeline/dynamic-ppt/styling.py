from pptx.enum.dml import MSO_FILL
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE

def KeepTextStyle(
    shape,
    new_text,
    new_size=None,
    paragraph_index=0
):
    """
    Replace text inside a shape without destroying existing font styling.
    """

    if not shape.has_text_frame:
        return

    tf = shape.text_frame
    tf.auto_size = None
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    if len(tf.paragraphs) <= paragraph_index:
        return

    paragraph = tf.paragraphs[paragraph_index]

    if not paragraph.runs:
        return

    # Keep original styling by reusing first run
    first_run = paragraph.runs[0]
    first_run.text = new_text

    if new_size is not None:
        first_run.font.size = Pt(new_size)

    # Clear remaining runs safely
    for run in paragraph.runs[1:]:
        run.text = ""

def UpdateShapeTextByName(slide,shape_name,new_text,new_size=None):
        """
        Find shape by name and update text and keeping style
        """
        for shape in slide.shapes:
            if shape.name == shape_name:
                KeepTextStyle(
                    shape,
                    new_text,
                    new_size=new_size)
                break