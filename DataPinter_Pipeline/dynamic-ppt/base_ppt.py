from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL
import os
import datetime
import pandas as pd
import sys
from copy import deepcopy
sys.path.append(r"C:\KRESNA\Tools-20251201T015117Z-1-001\Tools\DATA PIPELINE\DataPinter_Pipeline")
sys.path.append(r"C:\KRESNA\Tools-20251201T015117Z-1-001\Tools\DATA PIPELINE\DataPinter_Pipeline\dynamic-ppt")
from styling import KeepTextStyle

#PATH
ASSET_PATH = r"C:\KRESNA\Tools-20251201T015117Z-1-001\Tools\DATA PIPELINE\DataPinter_Pipeline\dynamic-ppt\asset\PPT_BASE.pptx"

def CreateBaseAndTitle(df, TEMPLATE_PATH, category, PATH, PPT_NAME):

    date = pd.to_datetime(df['query_date']).iloc[0]
    year = date.year
    month = date.strftime("%B")

    prs = Presentation(TEMPLATE_PATH)
    slide = prs.slides[0]

    for shape in slide.shapes:
        if shape.has_text_frame:
            if shape.name == "TextBox 3":
                KeepTextStyle(shape,new_text=f'{category} Product Market',
                              new_size = 100)

            elif shape.name == "TextBox 4":
                KeepTextStyle(shape, new_text=f'{month} {year}',
                              new_size = 30)
    prs.save(os.path.join(PATH, PPT_NAME))
    return prs
    

def AddCharts(prs,page_id,chart_img, size, position):
    """
    Function to add image chart into PPT. 
    Params:
    1.Slides:Unique number indexing the slide position relative to the whole PPT file
    2.chart_img : Image file name and directory
    2.size : Size of the image desired: Tuple type file > (width, height)
    3.position : Position of image inside the PPT Slides : Tuple Type with format (left, top)
    """
    if chart_img.lower().endswith((".png",".jpg",".jpeg")):
        width, height = size
        left, top = position

        #Get the page id, and put the pictures into slide
        slide = prs.slides[page_id]
        pic= slide.shapes.add_picture(chart_img, 
                                      left = left, 
                                      top = top, 
                                      width = width, 
                                      height = height)
        return pic
    else: 
        raise ValueError(" File is not a valid image (.png/.jpg/jpeg)")

def AddTextBox(prs,page_id, text,textfont,
               textsize,textcolor,
               size,position,
               bold=False,
               align=None):
    
    left, top = position
    width, height = size

    slide = prs.slides[page_id]
    txtbox = slide.shapes.add_textbox(left, top,width, height)

    tf = txtbox.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = text

    run = p.runs[0]
    run.font.name = textfont
    run.font.size = textsize
    run.font.color.rgb = textcolor

    return txtbox

def DuplicateSlide(prs,
                   source_index,
                   insert_after_index):
    """
    Duplicate the slide and insert it after specific slide pages

    Params:
    - prs
    - source_index: Page Index that will be copied into new slide
    - insert_after_index : Slide Position relative to the whole presentation slides
    """

    source_slide = prs.slides[source_index]

    #Create new slide
    blank_slide = source_slide.slide_layout
    new_slide = prs.slides.add_slide(blank_slide)

    for shape in source_slide.shapes:
        new_el = deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(
            new_el, 'p:extLst')
    slide_id_list = prs.slides._sldIdLst
    slides = list(slide_id_list)

    slide = slides[-1]
    slide_id_list.remove(slide)
    slide_id_list.insert(insert_after_index + 1, slide)

    return new_slide


